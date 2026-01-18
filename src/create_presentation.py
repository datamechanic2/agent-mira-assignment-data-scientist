"""
Generate PowerPoint presentation for the case study
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path


def add_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER


def add_bullets(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # bullets
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = f"    - {item[1]}"
            p.font.size = Pt(18)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)
        p.space_after = Pt(10)


def add_image(prs, title, img_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(0.75), Inches(0.9), width=Inches(8.5))

    if caption:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(8.5), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True


def add_table(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.2),
                                    Inches(9), Inches(0.4 * n_rows)).table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r + 1, c).text = str(val)


def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    base = Path(__file__).parent.parent
    eda = base / "eda_outputs"

    # slides
    add_title(prs, "Real Estate Price Prediction", "ML Case Study - Agent Mira")

    add_bullets(prs, "Summary", [
        "Goal: Predict house prices using ML",
        "Data: 247K transactions (2020-2024)",
        "Best model: LightGBM (R²=0.53, RMSE=$157K)",
        "Deployed as REST API with multi-core support",
        "Size and sale timing are top predictors"
    ])

    add_bullets(prs, "Dataset", [
        "247,172 property sales",
        "Features:",
        ("sub", "Location (4 cities)"),
        ("sub", "Size (800-4000 sqft)"),
        ("sub", "Bedrooms (1-5), Bathrooms (1-3)"),
        ("sub", "Year Built, Condition, Type"),
        "Target: Price ($26K - $2.2M)"
    ])

    add_bullets(prs, "Data Quality", [
        "Missing values handled:",
        ("sub", "Year Built: 5.1%"),
        ("sub", "Condition: 4.3%"),
        ("sub", "Bedrooms: 3.4%"),
        ("sub", "Price: 2.2%"),
        "Strategy: median for numeric, mode for categorical"
    ])

    if (eda / "price_distribution.png").exists():
        add_image(prs, "Price Distribution", eda / "price_distribution.png",
                  "Right-skewed | Mean: $466K | Median: $417K")

    if (eda / "categorical_analysis.png").exists():
        add_image(prs, "Categorical Analysis", eda / "categorical_analysis.png")

    if (eda / "correlation_matrix.png").exists():
        add_image(prs, "Correlations", eda / "correlation_matrix.png")

    add_bullets(prs, "Feature Engineering", [
        "Created 12 features:",
        ("sub", "property_age, total_rooms, bath_ratio"),
        ("sub", "size_cat (small/medium/large/xlarge)"),
        ("sub", "year_sold, month_sold, quarter"),
        ("sub", "is_new, decade"),
        "One-hot encoding for categoricals",
        "StandardScaler for numerics"
    ])

    add_table(prs, "Model Comparison",
        ["Model", "Val RMSE", "R²", "Time"],
        [
            ["LightGBM", "$158K", "0.53", "0.7s"],
            ["XGBoost", "$158K", "0.53", "0.3s"],
            ["GBM", "$158K", "0.53", "18.7s"],
            ["RandomForest", "$161K", "0.52", "3.2s"],
            ["Ridge", "$163K", "0.50", "0.01s"]
        ])

    add_bullets(prs, "Best Model: LightGBM", [
        "Test performance:",
        ("sub", "RMSE: $157,307"),
        ("sub", "MAE: $117,716"),
        ("sub", "R²: 0.5317"),
        "Tuned with GridSearchCV (5-fold CV)",
        "Params: lr=0.05, depth=4, n_est=200"
    ])

    add_table(prs, "Feature Importance",
        ["Rank", "Feature", "Score"],
        [
            ["1", "Size", "850"],
            ["2", "year_sold", "469"],
            ["3", "month_sold", "439"],
            ["4", "total_rooms", "359"],
            ["5", "property_age", "326"]
        ])

    add_bullets(prs, "Multi-Core Processing", [
        "Training: n_jobs for RF, XGB, LGBM",
        "GridSearchCV with parallel CV folds",
        "API: ThreadPoolExecutor for batch predictions",
        "Uvicorn with multiple workers"
    ])

    add_bullets(prs, "API Endpoints", [
        "POST /predict - single property",
        "POST /predict/batch - multiple properties",
        "GET /health - status check",
        "GET /model/info - model details",
        "Swagger docs at /docs"
    ])

    add_bullets(prs, "Project Files", [
        "src/eda.py - EDA and plots",
        "src/preprocessing.py - data pipeline",
        "src/model_training.py - model training",
        "api/app.py - FastAPI server",
        "main.py - orchestration",
        "models/ - saved model artifacts"
    ])

    add_title(prs, "Thank You", "Questions?")

    out = base / "Real_Estate_Price_Prediction_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    create_pptx()
